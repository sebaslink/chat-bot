#!/usr/bin/env python3
"""
Cargador independiente de datos para el chatbot
Permite cargar URLs y PDFs antes de usar el chat
"""

import json
import os
import re
from typing import List, Dict

class DataLoader:
    def __init__(self, data_file="knowledge_base.json"):
        self.data_file = data_file
        self.documents = []
        self.load_data()
    
    def load_data(self):
        """Carga datos existentes del archivo JSON"""
        if os.path.exists(self.data_file):
            try:
                with open(self.data_file, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    self.documents = data.get('documents', [])
                print(f"✅ Cargados {len(self.documents)} documentos existentes")
            except Exception as e:
                print(f"⚠️ Error cargando datos existentes: {e}")
                self.documents = []
    
    def save_data(self):
        """Guarda los datos en el archivo JSON"""
        try:
            data = {
                'documents': self.documents,
                'total_documents': len(self.documents)
            }
            with open(self.data_file, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            print(f"✅ Datos guardados: {len(self.documents)} documentos")
        except Exception as e:
            print(f"❌ Error guardando datos: {e}")
    
    def _chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Divide el texto en fragmentos más pequeños"""
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            
            # Intentar cortar en un punto lógico
            if end < len(text):
                last_period = chunk.rfind('.')
                last_newline = chunk.rfind('\n')
                cut_point = max(last_period, last_newline)
                
                if cut_point > start + chunk_size // 2:
                    chunk = chunk[:cut_point + 1]
                    end = start + cut_point + 1
            
            chunks.append(chunk.strip())
            start = end - overlap
        
        return [chunk for chunk in chunks if chunk.strip()]
    
    def add_text_document(self, content: str, source: str = "manual", title: str = ""):
        """Añade un documento de texto a la base de conocimientos"""
        if not content.strip():
            print("⚠️ El contenido no puede estar vacío")
            return False
        
        # Dividir en chunks si es muy largo
        chunks = self._chunk_text(content)
        
        for i, chunk in enumerate(chunks):
            if chunk.strip():
                self.documents.append({
                    'id': f"{source}_{len(self.documents)}",
                    'content': chunk,
                    'source': source,
                    'title': title,
                    'chunk_index': i,
                    'total_chunks': len(chunks)
                })
        
        self.save_data()
        print(f"✅ Añadido documento '{title}' con {len(chunks)} fragmentos")
        return True
    
    def load_pdf(self, pdf_path: str) -> bool:
        """Carga un archivo PDF con soporte para OCR"""
        try:
            import PyPDF2
            
            # Verificar que el archivo existe
            if not os.path.exists(pdf_path):
                print(f"❌ El archivo no existe: {pdf_path}")
                return False
            
            # Verificar que el archivo no está vacío
            if os.path.getsize(pdf_path) == 0:
                print(f"❌ El archivo está vacío: {pdf_path}")
                return False
            
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Verificar que el PDF tiene páginas
                if len(pdf_reader.pages) == 0:
                    print(f"❌ El PDF no tiene páginas: {pdf_path}")
                    return False
                
                text = ""
                
                # Procesar todas las páginas
                for page_num in range(len(pdf_reader.pages)):
                    try:
                        page = pdf_reader.pages[page_num]
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                    except Exception as e:
                        print(f"⚠️ Error procesando página {page_num + 1} de {pdf_path}: {str(e)}")
                        continue
                
                # Si no se pudo extraer texto, intentar con OCR
                if not text.strip():
                    print(f"⚠️ No se pudo extraer texto directamente. Intentando con OCR...")
                    text = self._extract_text_with_ocr(pdf_path)
                
                if text.strip():
                    filename = os.path.basename(pdf_path)
                    return self.add_text_document(text, "pdf", filename)
                else:
                    print(f"❌ No se pudo extraer texto del PDF: {pdf_path}")
                    return False
                    
        except ImportError:
            print("❌ PyPDF2 no está instalado. Instálalo con: pip install PyPDF2")
            return False
        except Exception as e:
            print(f"❌ Error procesando PDF {pdf_path}: {str(e)}")
            return False
    
    def _extract_text_with_ocr(self, pdf_path: str) -> str:
        """Extrae texto de PDF usando OCR"""
        try:
            import fitz  # PyMuPDF
            from PIL import Image
            import pytesseract
            import io
            
            print("🔍 Iniciando extracción con OCR...")
            
            # Abrir el PDF con PyMuPDF
            doc = fitz.open(pdf_path)
            text = ""
            
            for page_num in range(len(doc)):
                try:
                    page = doc[page_num]
                    # Convertir la página a imagen
                    pix = page.get_pixmap()
                    img_data = pix.tobytes("png")
                    img = Image.open(io.BytesIO(img_data))
                    
                    # Extraer texto con OCR
                    page_text = pytesseract.image_to_string(img, lang='spa+eng')
                    if page_text.strip():
                        text += page_text + "\n"
                        print(f"✅ Página {page_num + 1} procesada con OCR")
                    
                except Exception as e:
                    print(f"⚠️ Error en OCR página {page_num + 1}: {str(e)}")
                    continue
            
            doc.close()
            return text
            
        except ImportError as e:
            print(f"⚠️ Librerías OCR no disponibles: {e}")
            print("💡 Para procesar PDFs escaneados, instala: pip install PyMuPDF pytesseract pillow")
            return ""
        except Exception as e:
            print(f"⚠️ Error en OCR: {str(e)}")
            return ""
    
    def load_pdfs_from_directory(self, directory_path: str) -> int:
        """Carga todos los PDFs de un directorio"""
        if not os.path.exists(directory_path):
            print(f"❌ El directorio {directory_path} no existe")
            return 0
        
        loaded_count = 0
        for filename in os.listdir(directory_path):
            if filename.lower().endswith('.pdf'):
                pdf_path = os.path.join(directory_path, filename)
                print(f"📄 Procesando: {filename}")
                if self.load_pdf(pdf_path):
                    loaded_count += 1
        
        return loaded_count
    
    def load_url(self, url: str) -> bool:
        """Carga contenido de una URL con múltiples estrategias"""
        try:
            import requests
            from bs4 import BeautifulSoup
            
            # Múltiples User-Agents para evitar bloqueos
            user_agents = [
                'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/119.0.0.0 Safari/537.36',
                'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                'Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'
            ]
            
            # Headers más completos
            headers = {
                'User-Agent': user_agents[0],
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'es-ES,es;q=0.9,en;q=0.8',
                'Accept-Encoding': 'gzip, deflate, br',
                'DNT': '1',
                'Connection': 'keep-alive',
                'Upgrade-Insecure-Requests': '1',
            }
            
            # Intentar con diferentes User-Agents
            for i, user_agent in enumerate(user_agents):
                try:
                    headers['User-Agent'] = user_agent
                    print(f"🌐 Intentando acceder a {url} (intento {i+1}/{len(user_agents)})")
                    
                    response = requests.get(url, headers=headers, timeout=15, allow_redirects=True)
                    
                    # Si es 403, intentar con el siguiente User-Agent
                    if response.status_code == 403:
                        print(f"⚠️ Error 403 con User-Agent {i+1}, probando siguiente...")
                        continue
                    
                    response.raise_for_status()
                    
                    # Si llegamos aquí, la petición fue exitosa
                    print(f"✅ Acceso exitoso a {url}")
                    break
                    
                except requests.exceptions.RequestException as e:
                    if i == len(user_agents) - 1:  # Último intento
                        raise e
                    print(f"⚠️ Error con User-Agent {i+1}: {str(e)}, probando siguiente...")
                    continue
            
            # Procesar el contenido
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # Remover scripts, estilos y otros elementos no deseados
            for script in soup(["script", "style", "nav", "footer", "header", "aside"]):
                script.decompose()
            
            # Extraer texto
            text = soup.get_text()
            
            # Limpiar texto
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            # Limitar longitud
            if len(text) > 15000:
                text = text[:15000] + "..."
            
            if text.strip():
                # Obtener título
                title = soup.find('title')
                title_text = title.get_text().strip() if title else url
                
                print(f"✅ Contenido extraído: {len(text)} caracteres")
                return self.add_text_document(text, "web", title_text)
            else:
                print(f"❌ No se pudo extraer contenido de la URL: {url}")
                return False
                
        except ImportError:
            print("❌ requests o beautifulsoup4 no están instalados. Instálalos con: pip install requests beautifulsoup4")
            return False
        except requests.exceptions.RequestException as e:
            print(f"❌ Error de conexión con {url}: {str(e)}")
            return False
        except Exception as e:
            print(f"❌ Error procesando URL {url}: {str(e)}")
            return False
    
    def load_urls(self, urls: List[str]) -> int:
        """Carga contenido de múltiples URLs"""
        loaded_count = 0
        for url in urls:
            print(f"🌐 Procesando: {url}")
            if self.load_url(url):
                loaded_count += 1
        
        return loaded_count
    
    def get_stats(self) -> Dict:
        """Obtiene estadísticas de la base de conocimientos"""
        sources = {}
        for doc in self.documents:
            source = doc.get('source', 'unknown')
            sources[source] = sources.get(source, 0) + 1
        
        return {
            'total_documents': len(self.documents),
            'sources': sources,
            'data_file': self.data_file
        }
    
    def clear_data(self):
        """Limpia todos los datos"""
        self.documents = []
        if os.path.exists(self.data_file):
            os.remove(self.data_file)
        print("✅ Base de conocimientos limpiada")
    
    def load_image(self, image_path: str, title: str = None) -> bool:
        """Carga una imagen y extrae texto usando OCR"""
        try:
            from PIL import Image
            import pytesseract
            
            # Verificar que el archivo existe
            if not os.path.exists(image_path):
                print(f"❌ El archivo no existe: {image_path}")
                return False
            
            # Verificar que el archivo no está vacío
            if os.path.getsize(image_path) == 0:
                print(f"❌ El archivo está vacío: {image_path}")
                return False
            
            print(f"🔍 Procesando imagen: {image_path}")
            
            # Abrir la imagen
            image = Image.open(image_path)
            
            # Extraer texto con OCR
            text = pytesseract.image_to_string(image, lang='spa+eng')
            
            if text.strip():
                # Usar el título proporcionado o el nombre del archivo
                if not title:
                    title = os.path.basename(image_path)
                
                print(f"✅ Texto extraído: {len(text)} caracteres")
                return self.add_text_document(text, "image", title)
            else:
                print(f"❌ No se pudo extraer texto de la imagen: {image_path}")
                return False
                
        except ImportError as e:
            print(f"❌ Librerías OCR no disponibles: {e}")
            print("💡 Para procesar imágenes, instala: pip install pytesseract pillow")
            return False
        except Exception as e:
            print(f"❌ Error procesando imagen {image_path}: {str(e)}")
            return False
    
    def load_multiple_images(self, image_paths: List[str]) -> int:
        """Carga múltiples imágenes y extrae texto usando OCR"""
        success_count = 0
        
        for image_path in image_paths:
            if self.load_image(image_path):
                success_count += 1
        
        print(f"✅ Procesadas {success_count}/{len(image_paths)} imágenes exitosamente")
        return success_count
    
    def load_word_document(self, docx_path: str, title: str = None) -> bool:
        """Carga un documento Word (.docx) y extrae texto"""
        try:
            from docx import Document
            
            # Verificar que el archivo existe
            if not os.path.exists(docx_path):
                print(f"❌ El archivo no existe: {docx_path}")
                return False
            
            # Verificar que el archivo no está vacío
            if os.path.getsize(docx_path) == 0:
                print(f"❌ El archivo está vacío: {docx_path}")
                return False
            
            print(f"📄 Procesando documento Word: {docx_path}")
            
            # Abrir el documento
            doc = Document(docx_path)
            
            # Extraer texto de todos los párrafos
            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text.strip())
            
            # Extraer texto de tablas
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text.strip())
            
            text = '\n'.join(text_parts)
            
            if text.strip():
                # Usar el título proporcionado o el nombre del archivo
                if not title:
                    title = os.path.basename(docx_path)
                
                print(f"✅ Texto extraído: {len(text)} caracteres")
                return self.add_text_document(text, "word", title)
            else:
                print(f"❌ No se pudo extraer texto del documento Word: {docx_path}")
                return False
                
        except ImportError as e:
            print(f"❌ Librería python-docx no disponible: {e}")
            print("💡 Para procesar documentos Word, instala: pip install python-docx")
            return False
        except Exception as e:
            print(f"❌ Error procesando documento Word {docx_path}: {str(e)}")
            return False
    
    def load_powerpoint_document(self, pptx_path: str, title: str = None) -> bool:
        """Carga una presentación PowerPoint (.pptx) y extrae texto"""
        try:
            from pptx import Presentation
            
            # Verificar que el archivo existe
            if not os.path.exists(pptx_path):
                print(f"❌ El archivo no existe: {pptx_path}")
                return False
            
            # Verificar que el archivo no está vacío
            if os.path.getsize(pptx_path) == 0:
                print(f"❌ El archivo está vacío: {pptx_path}")
                return False
            
            print(f"📊 Procesando presentación PowerPoint: {pptx_path}")
            
            # Abrir la presentación
            prs = Presentation(pptx_path)
            
            # Extraer texto de todas las diapositivas
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    text_parts.append(f"Diapositiva {slide_num}:")
                    text_parts.extend(slide_text)
                    text_parts.append("")  # Línea en blanco entre diapositivas
            
            text = '\n'.join(text_parts)
            
            if text.strip():
                # Usar el título proporcionado o el nombre del archivo
                if not title:
                    title = os.path.basename(pptx_path)
                
                print(f"✅ Texto extraído: {len(text)} caracteres")
                return self.add_text_document(text, "powerpoint", title)
            else:
                print(f"❌ No se pudo extraer texto de la presentación PowerPoint: {pptx_path}")
                return False
                
        except ImportError as e:
            print(f"❌ Librería python-pptx no disponible: {e}")
            print("💡 Para procesar presentaciones PowerPoint, instala: pip install python-pptx")
            return False
        except Exception as e:
            print(f"❌ Error procesando presentación PowerPoint {pptx_path}: {str(e)}")
            return False
    
    def load_excel_document(self, xlsx_path: str, title: str = None) -> bool:
        """Carga una hoja de cálculo Excel (.xlsx) y extrae texto"""
        try:
            import pandas as pd
            
            # Verificar que el archivo existe
            if not os.path.exists(xlsx_path):
                print(f"❌ El archivo no existe: {xlsx_path}")
                return False
            
            # Verificar que el archivo no está vacío
            if os.path.getsize(xlsx_path) == 0:
                print(f"❌ El archivo está vacío: {xlsx_path}")
                return False
            
            print(f"📈 Procesando hoja de cálculo Excel: {xlsx_path}")
            
            # Leer todas las hojas del archivo Excel
            excel_file = pd.ExcelFile(xlsx_path)
            text_parts = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(xlsx_path, sheet_name=sheet_name)
                
                # Agregar nombre de la hoja
                text_parts.append(f"Hoja: {sheet_name}")
                
                # Convertir DataFrame a texto
                # Reemplazar NaN con cadenas vacías
                df = df.fillna('')
                
                # Agregar encabezados
                if not df.empty:
                    headers = ' | '.join([str(col) for col in df.columns])
                    text_parts.append(f"Encabezados: {headers}")
                    
                    # Agregar datos (máximo 100 filas para evitar texto muy largo)
                    for index, row in df.head(100).iterrows():
                        row_text = ' | '.join([str(cell) for cell in row.values if str(cell).strip()])
                        if row_text.strip():
                            text_parts.append(row_text)
                    
                    if len(df) > 100:
                        text_parts.append(f"... y {len(df) - 100} filas más")
                
                text_parts.append("")  # Línea en blanco entre hojas
            
            text = '\n'.join(text_parts)
            
            if text.strip():
                # Usar el título proporcionado o el nombre del archivo
                if not title:
                    title = os.path.basename(xlsx_path)
                
                print(f"✅ Texto extraído: {len(text)} caracteres")
                return self.add_text_document(text, "excel", title)
            else:
                print(f"❌ No se pudo extraer texto de la hoja de cálculo Excel: {xlsx_path}")
                return False
                
        except ImportError as e:
            print(f"❌ Librería pandas no disponible: {e}")
            print("💡 Para procesar hojas de cálculo Excel, instala: pip install pandas openpyxl")
            return False
        except Exception as e:
            print(f"❌ Error procesando hoja de cálculo Excel {xlsx_path}: {str(e)}")
            return False
    
    def load_document_by_type(self, file_path: str, title: str = None) -> bool:
        """Carga un documento basándose en su extensión"""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.pdf':
            return self.load_pdf(file_path)
        elif file_extension == '.docx':
            return self.load_word_document(file_path, title)
        elif file_extension == '.pptx':
            return self.load_powerpoint_document(file_path, title)
        elif file_extension in ['.xlsx', '.xls']:
            return self.load_excel_document(file_path, title)
        elif file_extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
            return self.load_image(file_path, title)
        else:
            print(f"❌ Tipo de archivo no soportado: {file_extension}")
            print("💡 Tipos soportados: PDF, DOCX, PPTX, XLSX, XLS, JPG, PNG, GIF, BMP, TIFF")
            return False

def main():
    """Función principal del cargador de datos"""
    print("📚 Cargador de Datos para Chatbot")
    print("=" * 50)
    
    loader = DataLoader()
    
    while True:
        print("\n" + "=" * 50)
        print("Opciones disponibles:")
        print("1. Cargar PDFs de un directorio")
        print("2. Cargar PDF individual")
        print("3. Cargar URLs")
        print("4. Añadir texto manual")
        print("5. Ver estadísticas")
        print("6. Limpiar base de datos")
        print("7. Salir")
        
        choice = input("\nSelecciona una opción (1-7): ").strip()
        
        if choice == "1":
            # Cargar PDFs de directorio
            directory = input("\nIngresa la ruta del directorio con PDFs: ").strip()
            if directory:
                count = loader.load_pdfs_from_directory(directory)
                print(f"✅ Cargados {count} PDFs exitosamente")
            else:
                print("⚠️ Por favor ingresa una ruta válida.")
        
        elif choice == "2":
            # Cargar PDF individual
            pdf_path = input("\nIngresa la ruta del archivo PDF: ").strip()
            if pdf_path and os.path.exists(pdf_path):
                if loader.load_pdf(pdf_path):
                    print("✅ PDF cargado exitosamente")
                else:
                    print("❌ Error cargando el PDF")
            else:
                print("❌ Archivo no encontrado.")
        
        elif choice == "3":
            # Cargar URLs
            print("\nIngresa las URLs (una por línea, termina con una línea vacía):")
            urls = []
            while True:
                url = input().strip()
                if not url:
                    break
                urls.append(url)
            
            if urls:
                count = loader.load_urls(urls)
                print(f"✅ Cargadas {count} URLs exitosamente")
            else:
                print("⚠️ No se ingresaron URLs válidas.")
        
        elif choice == "4":
            # Añadir texto manual
            print("\nAñadir información manual:")
            title = input("Título del documento: ").strip()
            source = input("Fuente (ej: 'manual', 'web', 'pdf'): ").strip() or "manual"
            
            print("\nIngresa el contenido (termina con una línea que contenga solo 'FIN'):")
            content_lines = []
            while True:
                line = input()
                if line.strip() == "FIN":
                    break
                content_lines.append(line)
            
            content = "\n".join(content_lines)
            if content.strip():
                if loader.add_text_document(content, source, title):
                    print("✅ Texto añadido exitosamente")
            else:
                print("⚠️ No se ingresó contenido válido.")
        
        elif choice == "5":
            # Ver estadísticas
            stats = loader.get_stats()
            print(f"\n📊 Estadísticas:")
            print(f"  - Total de documentos: {stats['total_documents']}")
            print(f"  - Archivo de datos: {stats['data_file']}")
            print(f"  - Por fuente:")
            for source, count in stats['sources'].items():
                print(f"    * {source}: {count} documentos")
        
        elif choice == "6":
            # Limpiar base de datos
            confirm = input("\n¿Estás seguro de que quieres limpiar la base de datos? (s/n): ").strip().lower()
            if confirm == 's':
                loader.clear_data()
            else:
                print("Operación cancelada.")
        
        elif choice == "7":
            # Salir
            print("\n👋 ¡Datos cargados! Ahora puedes usar el chatbot.")
            break
        
        else:
            print("⚠️ Opción no válida. Por favor selecciona 1-7.")

if __name__ == "__main__":
    main()
